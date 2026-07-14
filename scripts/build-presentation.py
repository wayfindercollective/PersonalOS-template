#!/usr/bin/env python3
"""
build-presentation.py — PersonalOS presentation builder

Build an HTML slide deck (presentable in any browser) and/or a PowerPoint
file from a JSON deck definition. Designed for local models (Qwen) and Claude:
one bash call, no extra services.

Usage:
  python3 scripts/build-presentation.py deck.json
  python3 scripts/build-presentation.py deck.json --format both --out workspace/uploads/my-talk
  python3 scripts/build-presentation.py --stdin --format html --out /tmp/deck

JSON schema (deck.json):
{
  "title": "Talk title",
  "author": "Your Name",
  "subtitle": "optional deck subtitle",
  "theme": "midnight" | "light" | "coral" | "forest" | "charcoal",
  "slides": [
    {"type": "title", "title": "...", "subtitle": "..."},
    {"type": "section", "title": "...", "subtitle": "..."},
    {"type": "content", "title": "...", "bullets": ["..."], "notes": "..."},
    {"type": "two-column", "title": "...",
     "left": {"heading": "...", "bullets": ["..."]},
     "right": {"heading": "...", "bullets": ["..."]}},
    {"type": "stats", "title": "...",
     "stats": [{"value": "397B", "label": "params"}, ...]},
    {"type": "quote", "quote": "...", "attribution": "..."},
    {"type": "code", "title": "...", "code": "...", "language": "bash"},
    {"type": "compare", "title": "...",
     "left_title": "Cloud", "left": ["..."],
     "right_title": "Local", "right": ["..."]}
  ]
}
"""

from __future__ import annotations

import argparse
import base64
import html
import json
import mimetypes
import re
import sys
from datetime import date
from pathlib import Path
from typing import Any

THEMES: dict[str, dict[str, str]] = {
    "midnight": {
        "bg": "#0b1220",
        "bg_alt": "#121a2b",
        "card": "#162033",
        "text": "#e8eefc",
        "muted": "#9aabc9",
        "accent": "#5b8cff",
        "accent2": "#22d3ee",
        "danger": "#f87171",
        "ok": "#34d399",
        "border": "#243149",
    },
    "charcoal": {
        "bg": "#141414",
        "bg_alt": "#1c1c1c",
        "card": "#242424",
        "text": "#f5f5f5",
        "muted": "#a3a3a3",
        "accent": "#f59e0b",
        "accent2": "#fbbf24",
        "danger": "#f87171",
        "ok": "#4ade80",
        "border": "#333",
    },
    "light": {
        "bg": "#f7f8fb",
        "bg_alt": "#ffffff",
        "card": "#ffffff",
        "text": "#111827",
        "muted": "#6b7280",
        "accent": "#2563eb",
        "accent2": "#0ea5e9",
        "danger": "#dc2626",
        "ok": "#059669",
        "border": "#e5e7eb",
    },
    "coral": {
        "bg": "#1a1014",
        "bg_alt": "#24151b",
        "card": "#2c1a22",
        "text": "#fce8ec",
        "muted": "#d4a5b0",
        "accent": "#f96167",
        "accent2": "#f9e795",
        "danger": "#fb7185",
        "ok": "#86efac",
        "border": "#3d2430",
    },
    "forest": {
        "bg": "#0d1510",
        "bg_alt": "#132018",
        "card": "#18261d",
        "text": "#e8f5e9",
        "muted": "#9bbb9f",
        "accent": "#4ade80",
        "accent2": "#a3e635",
        "danger": "#f87171",
        "ok": "#34d399",
        "border": "#24362a",
    },
}

# PPTX theme colors (hex without #)
PPTX_THEMES: dict[str, dict[str, str]] = {
    "midnight": {"bg": "0B1220", "text": "E8EEFC", "muted": "9AABC9", "accent": "5B8CFF", "card": "162033"},
    "charcoal": {"bg": "141414", "text": "F5F5F5", "muted": "A3A3A3", "accent": "F59E0B", "card": "242424"},
    "light": {"bg": "F7F8FB", "text": "111827", "muted": "6B7280", "accent": "2563EB", "card": "FFFFFF"},
    "coral": {"bg": "1A1014", "text": "FCE8EC", "muted": "D4A5B0", "accent": "F96167", "card": "2C1A22"},
    "forest": {"bg": "0D1510", "text": "E8F5E9", "muted": "9BBB9F", "accent": "4ADE80", "card": "18261D"},
}


ASSETS = Path("./skills/presentation/assets")

# Constrained creative packs (model picks enums; CSS is deterministic)
MOTIFS = ("orbs", "mesh", "grid", "bars", "aurora", "none")
VIBES = ("keynote", "product", "technical", "bold")
BRANDS = ("personalos", "wayfinder", "none")
LAYOUTS_CONTENT = ("list", "cards", "numbered", "spotlight")
LAYOUTS_TITLE = ("center", "left", "bold")

# Optional accent overrides layered on themes
ACCENTS = {
    "electric": ("#5b8cff", "#22d3ee"),
    "amber": ("#f59e0b", "#fbbf24"),
    "mint": ("#34d399", "#a3e635"),
    "rose": ("#f96167", "#f9e795"),
    "violet": ("#a78bfa", "#22d3ee"),
    "default": (None, None),
}


def slugify(s: str) -> str:
    s = s.strip().lower()
    s = re.sub(r"[^a-z0-9]+", "-", s)
    return s.strip("-")[:60] or "presentation"


def file_to_data_uri(p: Path) -> str | None:
    if not p.exists() or not p.is_file():
        return None
    mime = mimetypes.guess_type(str(p))[0] or "application/octet-stream"
    if p.suffix.lower() == ".svg":
        mime = "image/svg+xml"
    data = base64.b64encode(p.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{data}"


def resolve_brand(deck: dict[str, Any]) -> dict[str, Any]:
    """Resolve logo/wordmark into embeddable brand chrome."""
    brand = (deck.get("brand") or "personalos").lower().strip()
    if brand not in BRANDS and brand not in ("custom",):
        # allow unknown as custom label
        pass
    logo = deck.get("logo") or deck.get("logo_path") or ""
    logo_text = deck.get("logo_text") or deck.get("brand_text") or ""
    show = deck.get("show_logo")
    if show is None:
        show = brand != "none"

    src = None
    if logo:
        logo_s = str(logo).strip()
        if logo_s.startswith("http://") or logo_s.startswith("https://") or logo_s.startswith("data:"):
            src = logo_s
        else:
            src = file_to_data_uri(Path(logo_s).expanduser())
    elif brand == "wayfinder":
        src = file_to_data_uri(ASSETS / "personalos-mark.svg")  # generic mark; replace with your logo
        logo_text = logo_text or "PersonalOS"
    elif brand == "none":
        show = False
    else:  # personalos default
        src = file_to_data_uri(ASSETS / "personalos-mark.svg")
        logo_text = logo_text or "PersonalOS"

    if not logo_text and brand == "personalos":
        logo_text = "PersonalOS"
    if not logo_text and brand == "wayfinder":
        logo_text = "PersonalOS"

    return {
        "show": bool(show),
        "src": src,
        "text": logo_text or "",
        "brand": brand,
    }


def load_deck(path: Path | None, use_stdin: bool) -> dict[str, Any]:
    if use_stdin:
        raw = sys.stdin.read()
    elif path:
        raw = path.read_text(encoding="utf-8")
    else:
        raise SystemExit("Provide a deck JSON path or --stdin")
    data = json.loads(raw)
    if not isinstance(data, dict) or "slides" not in data:
        raise SystemExit("Deck JSON must be an object with a 'slides' array")
    if not isinstance(data["slides"], list) or not data["slides"]:
        raise SystemExit("Deck must have at least one slide")
    data.setdefault("title", "Presentation")
    data.setdefault("author", "")
    data.setdefault("subtitle", "")
    data.setdefault("theme", "midnight")
    if data["theme"] not in THEMES:
        data["theme"] = "midnight"
    data.setdefault("motif", "orbs")
    if data.get("motif") not in MOTIFS:
        data["motif"] = "orbs"
    data.setdefault("vibe", "keynote")
    if data.get("vibe") not in VIBES:
        data["vibe"] = "keynote"
    data.setdefault("brand", "personalos")
    data.setdefault("accent", "default")
    if data.get("accent") not in ACCENTS:
        data["accent"] = "default"
    return data


def esc(s: Any) -> str:
    return html.escape("" if s is None else str(s), quote=True)


def brand_chrome(brand: dict[str, Any], corner: str = "tl") -> str:
    if not brand.get("show"):
        return ""
    img = ""
    if brand.get("src"):
        img = f'<img class="brand-logo" src="{esc(brand["src"])}" alt=""/>'
    text = f'<span class="brand-text">{esc(brand.get("text") or "")}</span>' if brand.get("text") else ""
    return f'<div class="brand-chrome brand-{corner}">{img}{text}</div>'


def bullets_html(items: list[Any] | None, style: str = "list") -> str:
    if not items:
        return ""
    lis = []
    for i, item in enumerate(items):
        if isinstance(item, dict):
            title = esc(item.get("title") or item.get("heading") or "")
            body = esc(item.get("body") or item.get("text") or "")
            icon = esc(item.get("icon") or "")
            if title and body:
                content = f"<strong>{title}</strong> <span class='b-sub'>{body}</span>"
            else:
                content = title or body
        else:
            content = esc(item)
            icon = ""
        if style == "numbered":
            lis.append(f'<li class="b-num"><span class="num">{i+1:02d}</span><span class="b-body">{content}</span></li>')
        elif style == "cards":
            mark = f'<span class="card-icon">{icon or "◆"}</span>'
            lis.append(f'<li class="b-card">{mark}<div class="b-body">{content}</div></li>')
        else:
            lis.append(f"<li>{content}</li>")
    cls = {
        "list": "bullets",
        "numbered": "bullets bullets-numbered",
        "cards": "bullets bullets-cards",
        "spotlight": "bullets",
    }.get(style, "bullets")
    return f"<ul class='{cls}'>{''.join(lis)}</ul>"


def slide_shell(
    stype: str,
    idx: int,
    total: int,
    notes: str,
    inner: str,
    brand: dict[str, Any],
    extra_class: str = "",
    layout: str = "",
) -> str:
    notes_attr = f' data-notes="{esc(notes)}"' if notes else ""
    layout_attr = f' data-layout="{esc(layout)}"' if layout else ""
    return f"""
<section class="slide {stype}-slide {extra_class}" data-index="{idx}"{notes_attr}{layout_attr}>
  <div class="slide-frame slide-bg">
    {brand_chrome(brand)}
    <div class="accent-rail" aria-hidden="true"></div>
    <div class="inner">
      {inner}
    </div>
    <div class="slide-footer">
      <span class="foot-brand">{esc(brand.get("text") or "")}</span>
      <span class="slide-num">{idx + 1} / {total}</span>
    </div>
  </div>
</section>"""


def render_slide(slide: dict[str, Any], idx: int, total: int, brand: dict[str, Any]) -> str:
    stype = (slide.get("type") or "content").lower()
    notes = slide.get("notes") or ""
    layout = (slide.get("layout") or "").lower()

    if stype == "title":
        layout = layout if layout in LAYOUTS_TITLE else "center"
        eyebrow = slide.get("eyebrow") or ""
        inner = f"""
      <div class="title-block layout-{esc(layout)}">
        <div class="title-glow" aria-hidden="true"></div>
        {f'<p class="eyebrow">{esc(eyebrow)}</p>' if eyebrow else ''}
        <h1>{esc(slide.get("title") or "")}</h1>
        {f'<p class="subtitle">{esc(slide.get("subtitle") or "")}</p>' if slide.get("subtitle") else ''}
        <div class="title-rule" aria-hidden="true"></div>
      </div>"""
        return slide_shell("title", idx, total, notes, inner, brand, "is-hero", layout)

    if stype == "section":
        inner = f"""
      <div class="section-block">
        <p class="eyebrow">Section</p>
        <h1>{esc(slide.get("title") or "")}</h1>
        {f'<p class="subtitle">{esc(slide.get("subtitle") or "")}</p>' if slide.get("subtitle") else ''}
      </div>"""
        return slide_shell("section", idx, total, notes, inner, brand, "is-section")

    if stype == "quote":
        inner = f"""
      <div class="quote-block">
        <div class="quote-mark" aria-hidden="true">“</div>
        <blockquote>{esc(slide.get("quote") or "")}</blockquote>
        {f'<p class="attribution">— {esc(slide.get("attribution") or "")}</p>' if slide.get("attribution") else ''}
      </div>"""
        return slide_shell("quote", idx, total, notes, inner, brand)

    if stype == "stats":
        stats = slide.get("stats") or []
        cards = []
        for i, s in enumerate(stats[:6]):
            cards.append(
                f"<div class='stat-card delay-{i}'>"
                f"<div class='stat-value'>{esc(s.get('value', ''))}</div>"
                f"<div class='stat-label'>{esc(s.get('label', ''))}</div>"
                f"</div>"
            )
        inner = f"""
      <h2>{esc(slide.get("title") or "")}</h2>
      <div class="stats-grid cols-{min(len(cards), 4)}">{''.join(cards)}</div>"""
        return slide_shell("stats", idx, total, notes, inner, brand)

    if stype == "two-column":
        left = slide.get("left") or {}
        right = slide.get("right") or {}
        if isinstance(left, list):
            left = {"bullets": left}
        if isinstance(right, list):
            right = {"bullets": right}
        inner = f"""
      <h2>{esc(slide.get("title") or "")}</h2>
      <div class="cols">
        <div class="col card">
          <div class="card-kicker"></div>
          <h3>{esc(left.get("heading") or left.get("title") or "")}</h3>
          {bullets_html(left.get("bullets"))}
          {f'<p class="body-text">{esc(left.get("text") or "")}</p>' if left.get("text") else ''}
        </div>
        <div class="col card">
          <div class="card-kicker kicker-2"></div>
          <h3>{esc(right.get("heading") or right.get("title") or "")}</h3>
          {bullets_html(right.get("bullets"))}
          {f'<p class="body-text">{esc(right.get("text") or "")}</p>' if right.get("text") else ''}
        </div>
      </div>"""
        return slide_shell("content", idx, total, notes, inner, brand, "is-columns")

    if stype == "compare":
        left_items = slide.get("left") or []
        right_items = slide.get("right") or []
        if isinstance(left_items, dict):
            left_items = left_items.get("bullets") or []
        if isinstance(right_items, dict):
            right_items = right_items.get("bullets") or []
        inner = f"""
      <h2>{esc(slide.get("title") or "")}</h2>
      <div class="cols compare-cols">
        <div class="col card compare-left">
          <div class="card-kicker"></div>
          <h3>{esc(slide.get("left_title") or "A")}</h3>
          {bullets_html(left_items)}
        </div>
        <div class="compare-vs" aria-hidden="true">VS</div>
        <div class="col card compare-right">
          <div class="card-kicker kicker-2"></div>
          <h3>{esc(slide.get("right_title") or "B")}</h3>
          {bullets_html(right_items)}
        </div>
      </div>"""
        return slide_shell("content", idx, total, notes, inner, brand, "is-compare")

    if stype == "code":
        code = slide.get("code") or ""
        lang = slide.get("language") or ""
        inner = f"""
      <div class="code-head">
        <h2>{esc(slide.get("title") or "")}</h2>
        {f'<span class="lang-pill">{esc(lang)}</span>' if lang else ''}
      </div>
      <pre class="code-block"><code>{esc(code)}</code></pre>"""
        return slide_shell("code", idx, total, notes, inner, brand)

    # content
    layout = layout if layout in LAYOUTS_CONTENT else "list"
    bullets = slide.get("bullets") or slide.get("points") or []
    body = slide.get("body") or slide.get("text") or ""
    style = "cards" if layout == "cards" else ("numbered" if layout == "numbered" else "list")
    if layout == "spotlight" and body:
        inner = f"""
      <div class="spotlight">
        <h2>{esc(slide.get("title") or "")}</h2>
        <p class="spotlight-text">{esc(body)}</p>
        {bullets_html(bullets, "list")}
      </div>"""
    else:
        inner = f"""
      <h2>{esc(slide.get("title") or "")}</h2>
      {f'<p class="lead">{esc(body)}</p>' if body else ''}
      {bullets_html(bullets, style)}"""
    return slide_shell("content", idx, total, notes, inner, brand, f"layout-{layout}", layout)


def motif_css(motif: str) -> str:
    if motif == "none":
        return ".slide-bg { background: var(--bg); }"
    if motif == "grid":
        return """
  .slide-bg {
    background-color: var(--bg);
    background-image:
      linear-gradient(color-mix(in srgb, var(--border) 55%, transparent) 1px, transparent 1px),
      linear-gradient(90deg, color-mix(in srgb, var(--border) 55%, transparent) 1px, transparent 1px),
      radial-gradient(ellipse at 20% 0%, color-mix(in srgb, var(--accent) 16%, transparent), transparent 50%);
    background-size: 48px 48px, 48px 48px, auto;
  }"""
    if motif == "mesh":
        return """
  .slide-bg {
    background:
      radial-gradient(ellipse at 10% 20%, color-mix(in srgb, var(--accent) 28%, transparent), transparent 42%),
      radial-gradient(ellipse at 90% 10%, color-mix(in srgb, var(--accent2) 20%, transparent), transparent 40%),
      radial-gradient(ellipse at 70% 90%, color-mix(in srgb, var(--accent) 14%, transparent), transparent 45%),
      var(--bg);
  }"""
    if motif == "bars":
        return """
  .slide-bg {
    background:
      linear-gradient(105deg,
        color-mix(in srgb, var(--accent) 14%, transparent) 0%,
        transparent 28%,
        transparent 72%,
        color-mix(in srgb, var(--accent2) 12%, transparent) 100%),
      var(--bg);
  }
  .slide-bg::after {
    content: "";
    position: absolute; left: 0; top: 0; bottom: 0; width: 8px;
    background: linear-gradient(180deg, var(--accent), var(--accent2));
    opacity: 0.9;
  }"""
    if motif == "aurora":
        return """
  .slide-bg {
    background:
      radial-gradient(ellipse at 30% 0%, color-mix(in srgb, var(--accent) 35%, transparent), transparent 50%),
      radial-gradient(ellipse at 80% 40%, color-mix(in srgb, var(--accent2) 25%, transparent), transparent 45%),
      radial-gradient(ellipse at 40% 100%, color-mix(in srgb, var(--accent) 18%, transparent), transparent 50%),
      var(--bg);
  }"""
    # orbs (default)
    return """
  .slide-bg {
    background:
      radial-gradient(ellipse at 15% 10%, color-mix(in srgb, var(--accent) 22%, transparent), transparent 45%),
      radial-gradient(ellipse at 90% 80%, color-mix(in srgb, var(--accent2) 16%, transparent), transparent 42%),
      var(--bg);
  }"""


def vibe_css(vibe: str) -> str:
    if vibe == "bold":
        return """
  h1 { font-weight: 800; letter-spacing: -0.03em; }
  h2 { font-weight: 800; }
  .bullets { gap: 1rem; font-size: clamp(1.15rem, 2.2vw, 1.45rem); }
  .stat-value { font-size: clamp(2.4rem, 5vw, 3.6rem); }
"""
    if vibe == "technical":
        return """
  body { font-family: ui-sans-serif, "Segoe UI", system-ui, sans-serif; }
  h1, h2, h3 { font-family: ui-monospace, "SF Mono", Menlo, Consolas, monospace; letter-spacing: -0.02em; }
  .eyebrow { font-family: ui-monospace, Menlo, monospace; }
  .card { border-radius: 10px; }
"""
    if vibe == "product":
        return """
  .card { border-radius: 22px; box-shadow: 0 20px 50px color-mix(in srgb, #000 25%, transparent); }
  .stat-card { border-radius: 22px; }
  h1 { font-weight: 700; }
"""
    # keynote
    return """
  h1 { font-weight: 750; }
  .title-block h1 { max-width: 14ch; }
"""


def build_html(deck: dict[str, Any]) -> str:
    theme = dict(THEMES[deck["theme"]])
    accent_name = deck.get("accent") or "default"
    a1, a2 = ACCENTS.get(accent_name, ACCENTS["default"])
    if a1:
        theme["accent"] = a1
        theme["accent2"] = a2 or a1

    brand = resolve_brand(deck)
    motif = deck.get("motif") or "orbs"
    vibe = deck.get("vibe") or "keynote"
    slides = deck["slides"]
    total = len(slides)
    slides_html = "\n".join(render_slide(s, i, total, brand) for i, s in enumerate(slides))
    title = esc(deck["title"])
    author = esc(deck.get("author") or "")
    subtitle = esc(deck.get("subtitle") or "")
    today = date.today().isoformat()

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>{title}</title>
<style>
  :root {{
    --bg: {theme["bg"]};
    --bg-alt: {theme["bg_alt"]};
    --card: {theme["card"]};
    --text: {theme["text"]};
    --muted: {theme["muted"]};
    --accent: {theme["accent"]};
    --accent2: {theme["accent2"]};
    --border: {theme["border"]};
    --ok: {theme["ok"]};
    --danger: {theme["danger"]};
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html, body {{
    height: 100%;
    background: var(--bg);
    color: var(--text);
    font-family: "Segoe UI", system-ui, -apple-system, sans-serif;
    overflow: hidden;
  }}
  #deck {{ height: 100vh; width: 100vw; position: relative; }}
  .slide {{
    display: none;
    position: absolute;
    inset: 0;
  }}
  .slide.active {{ display: block; }}
  .slide-frame {{
    position: relative;
    height: 100%;
    width: 100%;
    padding: clamp(1.4rem, 3.5vw, 3rem) clamp(1.6rem, 4vw, 3.4rem);
    display: flex;
    flex-direction: column;
    justify-content: center;
    overflow: hidden;
  }}
  .slide-frame::before {{
    content: "";
    position: absolute; inset: 0; z-index: 0;
  }}
{motif_css(motif)}
  .brand-chrome {{
    position: absolute;
    top: clamp(1rem, 2.5vw, 1.6rem);
    left: clamp(1.2rem, 3vw, 2.2rem);
    display: flex;
    align-items: center;
    gap: 0.65rem;
    z-index: 5;
    opacity: 0.92;
  }}
  .brand-logo {{
    height: 36px; width: auto; max-width: 120px;
    object-fit: contain;
    border-radius: 8px;
  }}
  .brand-text {{
    font-size: 0.85rem;
    font-weight: 650;
    letter-spacing: 0.04em;
    color: var(--muted);
    text-transform: uppercase;
  }}
  .accent-rail {{
    position: absolute;
    left: 0; top: 18%; bottom: 18%;
    width: 4px;
    border-radius: 0 4px 4px 0;
    background: linear-gradient(180deg, var(--accent), var(--accent2));
    opacity: 0.85;
    z-index: 2;
  }}
  .title-slide .accent-rail, .section-slide .accent-rail {{ opacity: 0; }}
  .inner {{
    position: relative;
    z-index: 3;
    max-width: 1120px;
    width: 100%;
    margin: 0 auto;
  }}
  h1 {{
    font-size: clamp(2.5rem, 6.2vw, 4.6rem);
    line-height: 1.05;
    letter-spacing: -0.025em;
  }}
  h2 {{
    font-size: clamp(1.7rem, 3.6vw, 2.7rem);
    line-height: 1.12;
    margin-bottom: 1.25rem;
    letter-spacing: -0.02em;
  }}
  h3 {{
    font-size: 1.2rem;
    color: var(--accent);
    margin-bottom: 0.75rem;
    font-weight: 700;
  }}
  .subtitle, .lead {{
    margin-top: 1rem;
    font-size: clamp(1.1rem, 2.1vw, 1.4rem);
    color: var(--muted);
    max-width: 38ch;
    line-height: 1.4;
  }}
  .eyebrow {{
    text-transform: uppercase;
    letter-spacing: 0.16em;
    font-size: 0.8rem;
    color: var(--accent2);
    margin-bottom: 0.85rem;
    font-weight: 700;
  }}
  .title-block {{ position: relative; }}
  .title-block.layout-left {{ text-align: left; }}
  .title-block.layout-center {{ text-align: center; margin: 0 auto; }}
  .title-block.layout-center .subtitle {{ margin-left: auto; margin-right: auto; }}
  .title-block.layout-bold h1 {{
    font-size: clamp(2.8rem, 7vw, 5.2rem);
    background: linear-gradient(120deg, var(--text) 20%, var(--accent2));
    -webkit-background-clip: text; background-clip: text;
    color: transparent;
  }}
  .title-glow {{
    position: absolute; inset: -20% -10% auto -10%; height: 140%;
    background: radial-gradient(ellipse at 30% 40%, color-mix(in srgb, var(--accent) 30%, transparent), transparent 55%);
    z-index: -1; pointer-events: none;
  }}
  .title-rule {{
    margin-top: 1.4rem;
    width: 88px; height: 4px; border-radius: 4px;
    background: linear-gradient(90deg, var(--accent), var(--accent2));
  }}
  .title-block.layout-center .title-rule {{ margin-left: auto; margin-right: auto; }}
  .section-block h1 {{ max-width: 16ch; }}
  .bullets {{
    list-style: none;
    display: flex;
    flex-direction: column;
    gap: 0.85rem;
    font-size: clamp(1.05rem, 2vw, 1.35rem);
    line-height: 1.4;
  }}
  .bullets li {{
    position: relative;
    padding-left: 1.4rem;
  }}
  .bullets:not(.bullets-numbered):not(.bullets-cards) li::before {{
    content: "";
    position: absolute;
    left: 0; top: 0.55em;
    width: 0.55rem; height: 0.55rem;
    border-radius: 999px;
    background: linear-gradient(135deg, var(--accent), var(--accent2));
    box-shadow: 0 0 12px color-mix(in srgb, var(--accent) 50%, transparent);
  }}
  .bullets-numbered li {{
    display: flex; gap: 0.9rem; align-items: flex-start;
    padding-left: 0;
  }}
  .bullets-numbered .num {{
    flex: 0 0 auto;
    font-weight: 800;
    font-size: 0.95rem;
    color: var(--accent);
    background: color-mix(in srgb, var(--accent) 14%, transparent);
    border: 1px solid color-mix(in srgb, var(--accent) 35%, transparent);
    border-radius: 10px;
    padding: 0.25rem 0.5rem;
    min-width: 2.2rem;
    text-align: center;
  }}
  .bullets-cards {{
    display: grid !important;
    grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
    gap: 0.9rem !important;
  }}
  .bullets-cards li {{
    padding: 1rem 1.1rem 1rem 1rem;
    background: color-mix(in srgb, var(--card) 92%, transparent);
    border: 1px solid var(--border);
    border-radius: 16px;
    display: flex; gap: 0.75rem; align-items: flex-start;
  }}
  .bullets-cards li::before {{ display: none; }}
  .card-icon {{ color: var(--accent2); font-size: 0.9rem; margin-top: 0.2rem; }}
  .b-sub {{ color: var(--muted); font-weight: 400; }}
  .cols {{
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 1.2rem;
    align-items: stretch;
  }}
  .compare-cols {{ grid-template-columns: 1fr auto 1fr; gap: 0.9rem; align-items: center; }}
  .compare-vs {{
    font-size: 0.75rem; font-weight: 800; letter-spacing: 0.12em;
    color: var(--muted);
    border: 1px solid var(--border);
    border-radius: 999px;
    padding: 0.45rem 0.6rem;
    background: var(--bg-alt);
  }}
  .card {{
    background: color-mix(in srgb, var(--card) 92%, transparent);
    border: 1px solid var(--border);
    border-radius: 18px;
    padding: 1.25rem 1.35rem 1.3rem;
    backdrop-filter: blur(8px);
    position: relative;
    overflow: hidden;
  }}
  .card-kicker {{
    position: absolute; left: 0; top: 0; right: 0; height: 3px;
    background: linear-gradient(90deg, var(--accent), transparent);
  }}
  .card-kicker.kicker-2 {{
    background: linear-gradient(90deg, var(--accent2), transparent);
  }}
  .stats-grid {{
    display: grid;
    gap: 1rem;
    margin-top: 0.4rem;
  }}
  .stats-grid.cols-1 {{ grid-template-columns: 1fr; }}
  .stats-grid.cols-2 {{ grid-template-columns: repeat(2, 1fr); }}
  .stats-grid.cols-3 {{ grid-template-columns: repeat(3, 1fr); }}
  .stats-grid.cols-4 {{ grid-template-columns: repeat(4, 1fr); }}
  .stat-card {{
    background: color-mix(in srgb, var(--card) 94%, transparent);
    border: 1px solid var(--border);
    border-radius: 18px;
    padding: 1.5rem 1.1rem;
    text-align: center;
    position: relative;
    overflow: hidden;
  }}
  .stat-card::after {{
    content: "";
    position: absolute; inset: auto -20% -40% -20%; height: 70%;
    background: radial-gradient(ellipse at 50% 0%, color-mix(in srgb, var(--accent) 22%, transparent), transparent 70%);
    pointer-events: none;
  }}
  .stat-value {{
    font-size: clamp(2rem, 4.2vw, 3.1rem);
    font-weight: 800;
    color: var(--accent);
    line-height: 1.05;
    position: relative; z-index: 1;
  }}
  .stat-label {{
    margin-top: 0.5rem;
    color: var(--muted);
    font-size: 0.95rem;
    position: relative; z-index: 1;
  }}
  .quote-block {{ max-width: 18em; position: relative; }}
  .quote-mark {{
    font-size: 5rem; line-height: 0.8; color: var(--accent);
    opacity: 0.35; font-family: Georgia, serif;
    margin-bottom: 0.4rem;
  }}
  .quote-slide blockquote {{
    font-size: clamp(1.55rem, 3.3vw, 2.5rem);
    line-height: 1.28;
    font-weight: 600;
  }}
  .attribution {{
    margin-top: 1.4rem;
    color: var(--muted);
    font-size: 1.05rem;
  }}
  .code-head {{
    display: flex; align-items: baseline; justify-content: space-between; gap: 1rem;
  }}
  .lang-pill {{
    font-size: 0.75rem; font-weight: 700; letter-spacing: 0.06em;
    text-transform: uppercase;
    color: var(--accent2);
    border: 1px solid var(--border);
    border-radius: 999px;
    padding: 0.3rem 0.7rem;
    background: color-mix(in srgb, var(--card) 80%, transparent);
  }}
  .code-block {{
    background: #070b14;
    border: 1px solid var(--border);
    border-radius: 14px;
    padding: 1.15rem 1.25rem;
    overflow: auto;
    max-height: 58vh;
    font-family: ui-monospace, "SF Mono", Menlo, Consolas, monospace;
    font-size: clamp(0.85rem, 1.45vw, 1.05rem);
    line-height: 1.5;
    color: #d7e3ff;
    white-space: pre-wrap;
    box-shadow: inset 0 0 0 1px color-mix(in srgb, var(--accent) 10%, transparent);
  }}
  .spotlight-text {{
    font-size: clamp(1.5rem, 3vw, 2.2rem);
    font-weight: 650;
    line-height: 1.25;
    max-width: 22ch;
    margin: 0.4rem 0 1.2rem;
    background: linear-gradient(120deg, var(--text), color-mix(in srgb, var(--accent2) 70%, var(--text)));
    -webkit-background-clip: text; background-clip: text; color: transparent;
  }}
  .slide-footer {{
    position: absolute;
    left: clamp(1.2rem, 3vw, 2.2rem);
    right: clamp(1.2rem, 3vw, 2.2rem);
    bottom: 0.95rem;
    display: flex;
    justify-content: space-between;
    align-items: center;
    z-index: 5;
    color: var(--muted);
    font-size: 0.78rem;
    opacity: 0.75;
  }}
  .foot-brand {{ letter-spacing: 0.04em; text-transform: uppercase; font-weight: 600; }}
  #progress {{
    position: fixed; left: 0; top: 0; height: 3px; width: 0%;
    background: linear-gradient(90deg, var(--accent), var(--accent2));
    z-index: 50; transition: width 0.2s ease;
  }}
  #hud {{
    position: fixed; left: 1rem; bottom: 0.9rem;
    display: flex; gap: 0.5rem; z-index: 40;
    opacity: 0.35; transition: opacity 0.2s;
  }}
  #hud:hover {{ opacity: 1; }}
  #hud button {{
    background: var(--card); color: var(--text);
    border: 1px solid var(--border); border-radius: 999px;
    padding: 0.4rem 0.8rem; font-size: 0.8rem; cursor: pointer;
  }}
  #help {{
    position: fixed; right: 1rem; top: 1rem;
    color: var(--muted); font-size: 0.78rem; opacity: 0.45; z-index: 40;
  }}
  .body-text:empty {{ display: none; }}
{vibe_css(vibe)}
  {motif_css(motif)}
  @media (max-width: 900px) {{
    .cols, .compare-cols {{ grid-template-columns: 1fr; }}
    .compare-vs {{ display: none; }}
    .stats-grid.cols-3, .stats-grid.cols-4 {{ grid-template-columns: repeat(2, 1fr); }}
  }}
  @media print {{
    html, body {{ overflow: visible; height: auto; }}
    .slide {{ display: block !important; position: relative; page-break-after: always; min-height: 100vh; }}
    #hud, #help, #progress {{ display: none; }}
  }}
</style>
</head>
<body class="vibe-{esc(vibe)} motif-{esc(motif)} brand-{esc(brand.get('brand') or 'personalos')}">
<div id="progress"></div>
<div id="help">← → space · F fullscreen · N notes</div>
<div id="hud">
  <button type="button" id="prevBtn" title="Previous">←</button>
  <button type="button" id="nextBtn" title="Next">→</button>
  <button type="button" id="fullBtn" title="Fullscreen">Fullscreen</button>
</div>
<main id="deck" aria-label="{title}">
{slides_html}
</main>
<script>
(function () {{
  const slides = Array.from(document.querySelectorAll('.slide'));
  let i = 0;
  const progress = document.getElementById('progress');
  function show(n) {{
    if (!slides.length) return;
    i = Math.max(0, Math.min(slides.length - 1, n));
    slides.forEach((s, idx) => s.classList.toggle('active', idx === i));
    progress.style.width = ((i + 1) / slides.length * 100) + '%';
    location.hash = String(i + 1);
  }}
  function next() {{ show(i + 1); }}
  function prev() {{ show(i - 1); }}
  document.getElementById('nextBtn').onclick = next;
  document.getElementById('prevBtn').onclick = prev;
  document.getElementById('fullBtn').onclick = () => {{
    if (!document.fullscreenElement) document.documentElement.requestFullscreen().catch(() => {{}});
    else document.exitFullscreen().catch(() => {{}});
  }};
  window.addEventListener('keydown', (e) => {{
    if (['ArrowRight', 'ArrowDown', 'PageDown', ' '].includes(e.key)) {{ e.preventDefault(); next(); }}
    if (['ArrowLeft', 'ArrowUp', 'PageUp', 'Backspace'].includes(e.key)) {{ e.preventDefault(); prev(); }}
    if (e.key === 'Home') {{ e.preventDefault(); show(0); }}
    if (e.key === 'End') {{ e.preventDefault(); show(slides.length - 1); }}
    if (e.key === 'f' || e.key === 'F') document.getElementById('fullBtn').click();
    if (e.key === 'n' || e.key === 'N') {{
      const notes = slides[i].dataset.notes;
      if (notes) alert('Speaker notes:\\n\\n' + notes);
    }}
  }});
  let touchX = null;
  window.addEventListener('touchstart', (e) => {{ touchX = e.changedTouches[0].screenX; }}, {{ passive: true }});
  window.addEventListener('touchend', (e) => {{
    if (touchX == null) return;
    const dx = e.changedTouches[0].screenX - touchX;
    if (Math.abs(dx) > 40) {{ if (dx < 0) next(); else prev(); }}
    touchX = null;
  }}, {{ passive: true }});
  const hash = parseInt((location.hash || '').replace('#', ''), 10);
  show(Number.isFinite(hash) && hash > 0 ? hash - 1 : 0);
  document.getElementById('deck').addEventListener('click', (e) => {{
    if (e.target.closest('#hud') || e.target.closest('a,button')) return;
    if (e.clientX >= window.innerWidth / 2) next(); else prev();
  }});
}})();
</script>
<!-- PersonalOS presentation · {today} · {author} · motif={esc(motif)} vibe={esc(vibe)} brand={esc(brand.get('brand') or '')} -->
</body>
</html>
"""


def build_pptx(deck: dict[str, Any], out_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError as e:
        raise SystemExit(
            "python-pptx is required for PPTX output. Install with: "
            "python3 -m pip install --user python-pptx"
        ) from e

    colors = PPTX_THEMES.get(deck["theme"], PPTX_THEMES["midnight"])
    bg = RGBColor.from_string(colors["bg"])
    text_c = RGBColor.from_string(colors["text"])
    muted_c = RGBColor.from_string(colors["muted"])
    accent_c = RGBColor.from_string(colors["accent"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_run(run, size: int, color: RGBColor, bold: bool = False) -> None:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"

    def fill_bg(slide) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()
        # send to back
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def add_title(slide, text: str, top: float = 0.7, size: int = 40) -> None:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(1.2))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        set_run(run, size, text_c, bold=True)

    def add_bullets(slide, items: list[Any], top: float = 2.1, left: float = 0.9, width: float = 11.5) -> None:
        if not items:
            return
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.level = 0
            p.space_after = Pt(10)
            run = p.add_run()
            if isinstance(item, dict):
                t = item.get("title") or item.get("heading") or ""
                b = item.get("body") or item.get("text") or ""
                run.text = f"• {t}: {b}" if t and b else f"• {t or b}"
            else:
                run.text = f"• {item}"
            set_run(run, 20, text_c)

    for slide_data in deck["slides"]:
        slide = prs.slides.add_slide(blank)
        fill_bg(slide)
        stype = (slide_data.get("type") or "content").lower()

        if stype == "title":
            add_title(slide, slide_data.get("title") or "", top=2.4, size=48)
            if slide_data.get("subtitle"):
                box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(1))
                p = box.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = slide_data["subtitle"]
                set_run(run, 22, muted_c)

        elif stype == "section":
            add_title(slide, slide_data.get("title") or "", top=2.6, size=44)
            if slide_data.get("subtitle"):
                box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(1))
                p = box.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = slide_data["subtitle"]
                set_run(run, 20, accent_c)

        elif stype == "quote":
            box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.8), Inches(3))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"“{slide_data.get('quote') or ''}”"
            set_run(run, 28, text_c, bold=True)
            if slide_data.get("attribution"):
                box2 = slide.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.8), Inches(0.6))
                p2 = box2.text_frame.paragraphs[0]
                run2 = p2.add_run()
                run2.text = f"— {slide_data['attribution']}"
                set_run(run2, 18, muted_c)

        elif stype == "stats":
            add_title(slide, slide_data.get("title") or "", top=0.6, size=34)
            stats = (slide_data.get("stats") or [])[:4]
            if stats:
                gap = 0.35
                card_w = (11.7 - gap * (len(stats) - 1)) / len(stats)
                for i, s in enumerate(stats):
                    left = 0.8 + i * (card_w + gap)
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(left),
                        Inches(2.4),
                        Inches(card_w),
                        Inches(3.0),
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(colors["card"])
                    shape.line.color.rgb = RGBColor.from_string(colors.get("border", colors["card"]))
                    tbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.0), Inches(card_w - 0.4), Inches(1.8))
                    tf = tbox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = str(s.get("value", ""))
                    set_run(run, 32, accent_c, bold=True)
                    p2 = tf.add_paragraph()
                    p2.alignment = PP_ALIGN.CENTER
                    run2 = p2.add_run()
                    run2.text = str(s.get("label", ""))
                    set_run(run2, 14, muted_c)

        elif stype in ("two-column", "compare"):
            add_title(slide, slide_data.get("title") or "", top=0.5, size=32)
            if stype == "compare":
                left_title = slide_data.get("left_title") or "A"
                right_title = slide_data.get("right_title") or "B"
                left_items = slide_data.get("left") or []
                right_items = slide_data.get("right") or []
            else:
                left = slide_data.get("left") or {}
                right = slide_data.get("right") or {}
                if isinstance(left, list):
                    left = {"bullets": left}
                if isinstance(right, list):
                    right = {"bullets": right}
                left_title = left.get("heading") or left.get("title") or ""
                right_title = right.get("heading") or right.get("title") or ""
                left_items = left.get("bullets") or []
                right_items = right.get("bullets") or []

            lhead = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5))
            r = lhead.text_frame.paragraphs[0].add_run()
            r.text = left_title
            set_run(r, 18, accent_c, bold=True)
            rhead = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5))
            r2 = rhead.text_frame.paragraphs[0].add_run()
            r2.text = right_title
            set_run(r2, 18, accent_c, bold=True)
            add_bullets(slide, left_items, top=2.3, left=0.8, width=5.5)
            add_bullets(slide, right_items, top=2.3, left=7.0, width=5.5)

        elif stype == "code":
            add_title(slide, slide_data.get("title") or "", top=0.5, size=32)
            box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = slide_data.get("code") or ""
            set_run(run, 14, text_c)
            run.font.name = "Consolas"

        else:
            add_title(slide, slide_data.get("title") or "", top=0.6, size=34)
            body = slide_data.get("body") or slide_data.get("text")
            top = 1.9
            if body:
                box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.0))
                p = box.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = body
                set_run(run, 18, muted_c)
                top = 2.8
            add_bullets(slide, slide_data.get("bullets") or slide_data.get("points") or [], top=top)

        # page number
        num = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.3))
        p = num.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(len(prs.slides))
        set_run(run, 10, muted_c)

    prs.save(str(out_path))


PRESENTATIONS_DIR = Path("./workspace/presentations")
UPLOADS_DIR = Path("./workspace/uploads")


def funnel_host() -> str:
    """Resolve this machine's Tailscale MagicDNS name for Funnel URLs."""
    import subprocess

    try:
        raw = subprocess.check_output(
            ["tailscale", "status", "--json"],
            timeout=5,
            stderr=subprocess.DEVNULL,
        )
        data = json.loads(raw)
        name = (data.get("Self") or {}).get("DNSName") or ""
        name = name.rstrip(".")
        if name:
            return name
    except Exception:
        pass
    return "your-host.tailnet.ts.net"


def publish_html(html_path: Path) -> str | None:
    """Public open URL via Tailscale Funnel (self-hosted). No third-party host.

    Deck is mirrored to workspace/presentations/ and served by personalos's
    presentation edge at /presentations/* with correct text/html headers.

    Requires Funnel to target http://127.0.0.1:8787 once:
      bash scripts/enable-presentation-funnel.sh
    """
    # Always mirror first (publish is just the public URL for that mirror)
    try:
        mirror_to_presentations(html_path)
    except OSError as e:
        print(f"publish warn: mirror failed: {e}", file=sys.stderr)
        return None

    host = funnel_host()
    url = f"https://{host}/presentations/{html_path.name}"
    return url


def write_present_url_sidecar(html_path: Path, url: str) -> Path:
    """Write a sidecar the Telegram bot reads to send a tappable open link."""
    sidecar = Path(str(html_path) + ".present-url")
    sidecar.write_text(url.strip() + "\n", encoding="utf-8")
    # Also under uploads/ if html is elsewhere, so auto-send scan can find it
    if UPLOADS_DIR not in html_path.parents and html_path.parent != UPLOADS_DIR:
        mirror = UPLOADS_DIR / (html_path.name + ".present-url")
        try:
            mirror.write_text(url.strip() + "\n", encoding="utf-8")
        except OSError:
            pass
    return sidecar


def mirror_to_presentations(html_path: Path) -> Path:
    """Copy HTML into the local presentations web root (Tailscale/LAN server)."""
    PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)
    dest = PRESENTATIONS_DIR / html_path.name
    dest.write_bytes(html_path.read_bytes())
    return dest


def main() -> int:
    parser = argparse.ArgumentParser(description="Build PersonalOS presentation decks")
    parser.add_argument("deck", nargs="?", help="Path to deck JSON")
    parser.add_argument("--stdin", action="store_true", help="Read deck JSON from stdin")
    parser.add_argument(
        "--format",
        choices=["html", "pptx", "both"],
        default="both",
        help="Output format (default: both)",
    )
    parser.add_argument(
        "--out",
        help="Output path prefix (without extension). Default: workspace/uploads/<slug>",
    )
    parser.add_argument("--open-hint", action="store_true", help="Print how to present the HTML deck")
    parser.add_argument(
        "--publish",
        action="store_true",
        default=True,
        help="Emit a Funnel public URL (https://<host>/presentations/...) (default: on)",
    )
    parser.add_argument(
        "--no-publish",
        action="store_true",
        help="Skip public URL (files only)",
    )
    args = parser.parse_args()
    do_publish = args.publish and not args.no_publish

    deck_path = Path(args.deck) if args.deck else None
    deck = load_deck(deck_path, args.stdin)

    if args.out:
        out_prefix = Path(args.out)
    else:
        UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
        out_prefix = UPLOADS_DIR / slugify(deck["title"])

    out_prefix.parent.mkdir(parents=True, exist_ok=True)
    written: list[str] = []
    public_url: str | None = None

    if args.format in ("html", "both"):
        html_path = out_prefix.with_suffix(".html")
        html_path.write_text(build_html(deck), encoding="utf-8")
        written.append(str(html_path))

        # Always mirror into presentations/ for the local open-anywhere-on-tailnet server
        try:
            mirrored = mirror_to_presentations(html_path)
            written.append(f"local-mirror: {mirrored}")
        except OSError as e:
            print(f"mirror warn: {e}", file=sys.stderr)

        if do_publish:
            public_url = publish_html(html_path)
            if public_url:
                write_present_url_sidecar(html_path, public_url)
                # Also write sidecar next to uploads copy if different
                uploads_html = UPLOADS_DIR / html_path.name
                if uploads_html != html_path and uploads_html.exists():
                    write_present_url_sidecar(uploads_html, public_url)
                elif html_path.parent == UPLOADS_DIR:
                    write_present_url_sidecar(html_path, public_url)

    if args.format in ("pptx", "both"):
        pptx_path = out_prefix.with_suffix(".pptx")
        build_pptx(deck, pptx_path)
        written.append(str(pptx_path))

    print("OK")
    for p in written:
        print(p)

    if public_url:
        print(f"PRESENT_URL={public_url}")
        print(f"Open on any device: {public_url}")

    if args.open_hint and any(".html" in p for p in written):
        html_p = next(p for p in written if p.endswith(".html"))
        if public_url:
            print(
                f"Present: tap {public_url} on any phone/laptop — opens in the browser. "
                "Arrow keys / space to advance."
            )
        else:
            print(
                f"Present: open {html_p} in a browser, press F for fullscreen, "
                "arrow keys / space to advance. (Public publish failed — file only.)"
            )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
